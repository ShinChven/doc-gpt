from pathlib import Path
import PyPDF2
import click
from docx import Document
from pptx import Presentation
import re
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse

from doc_gpt.config import get_config
from .ai_client import AIClient

def is_valid_url(url):
    url_pattern = re.compile(
        r'^https?://'  # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+[A-Z]{2,6}\.?|'  # domain...
        r'localhost|'  # localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})'  # ...or ip
        r'(?::\d+)?'  # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    return url_pattern.match(url) is not None

def scrape_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        # Extract text from p, h1-h6, li tags
        text = ' '.join([tag.get_text() for tag in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li'])])
        return text
    except requests.RequestException as e:
        print(f"Error scraping URL {url}: {str(e)}")
        return ""

def process_input(input_path):
    if is_valid_url(input_path):
        return scrape_url(input_path)
    
    input_path = Path(input_path)
    
    if not input_path.exists():
        print(f"Error: {input_path} does not exist.")
        return ""
    
    if input_path.is_dir():
        content = ""
        for file in input_path.glob('*'):
            content += process_file(file) + "\n\n"
        return content.strip()
    else:
        return process_file(input_path)

def process_file(file_path):
    ext = file_path.suffix.lower()
    
    try:
        if ext in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.pdf':
            return process_pdf(file_path)
        elif ext == '.docx':
            return process_docx(file_path)
        elif ext == '.pptx':
            return process_pptx(file_path)
        else:
            print(f"Warning: Unsupported file type {ext} for {file_path}")
            return ""
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}")
        return ""

def process_pdf(file_path):
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        return "\n".join(page.extract_text() for page in reader.pages)

def process_docx(file_path):
    doc = Document(file_path)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)

def process_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content.append(shape.text)
    return "\n".join(text_content)

def url_to_valid_filename(url):
    # Remove the protocol (http:// or https://)
    url = re.sub(r'^https?://', '', url)
    # Replace invalid filename characters with underscores
    url = re.sub(r'[\\/*?:"<>|=]', '_', url)
    # Replace slashes with hyphens
    url = url.replace('/', '-')
    # Limit the filename length (adjust as needed)
    max_length = 200
    if len(url) > max_length:
        url = url[:max_length]
    return f"{url}.doc-gpt.md"

def write_output(content, output, input_file):
    if is_valid_url(input_file):
        input_path = Path(url_to_valid_filename(input_file))
    else:
        input_path = Path(input_file)
    
    output = Path(output or Path.cwd())
    
    # Create the output directory if it's a new path with no extension
    if not output.exists() and output.suffix == "":
        output.mkdir(parents=True, exist_ok=True)

    # If output is a directory, set the output to a file within that directory
    if output.is_dir():
        output = output / input_path.name
        
    # Create the output file if it doesn't exist
    if not output.exists():
        output.touch()

    # Write or append content to the output file
    with open(str(output), 'a', encoding='utf-8') as f:
        if output.stat().st_size > 0:
            file_divider = f"\n------\n\n"
            f.write(file_divider)
        f.write(content + "\n")
    
    print(f"Output written to {str(output)}")

def process_task(input_file, output_file, model_alias, prompt_file, instructions_file, write_prompt=False, max_tokens=None):
    try:
        client = AIClient(get_config())
        input_text = process_input(input_file)

        if prompt_file:
            prompt = process_input(prompt_file)
        else:
            default_prompt_file = Path.cwd() / 'prompt.md'
            if default_prompt_file.exists():
                prompt = process_input(str(default_prompt_file))
            else:
                prompt = click.prompt("Enter your prompt", type=str)

        if not prompt.strip():
            raise click.UsageError("Prompt cannot be empty")

        instructions = ""
        if instructions_file:
            instructions = process_input(instructions_file)
        else:
            default_instructions_file = Path.cwd() / 'instructions.md'
            if default_instructions_file.exists():
                instructions = process_input(str(default_instructions_file))

        messages = []
        if instructions:
            messages.append({"role": "system", "content": instructions})

        message = """
<ProvidedDocument>
[document]
</ProvidedDocument>

[prompt]
""".replace("[document]", input_text).replace("[prompt]", prompt)
        messages.append({"role": "user", "content": message})
        print(f'Processing: "{input_file}"')
        response = client.request(messages, model_alias, max_tokens)

        # Decide what to write based on the write_prompt flag
        if write_prompt:
            formatted_prompt = "\n".join(
                f"**Role:** {m['role']}\n**Content:**\n{m['content']}\n"
                for m in messages
            )
            combined_content = f"# Prompt:\n{formatted_prompt}\n\n# Response:\n{response}\n"
            write_output(combined_content, output_file, input_file)
        else:
            write_output(response, output_file, input_file)

        click.echo("Content generation completed successfully.")
    except click.UsageError as e:
        click.echo(f"Usage error: {str(e)}", err=True)
    except click.ClickException as e:
        click.echo(str(e), err=True)
    except Exception as e:
        click.echo(f"An error occurred: {str(e)}", err=True)
